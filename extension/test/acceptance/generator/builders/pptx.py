"""PPTX builders — presentation concealment vectors from the Threat Model.

Vectors (CLAUDE.md PPTX doctrine): speaker notes (text hidden in the notes pane,
never shown on the slide) and off-slide shapes (negative EMU coordinates put the
textbox outside the slide canvas).
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Emu

from payloads import BODY_EN


def _prs() -> tuple[Presentation, object]:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title-only
    slide.shapes.title.text = BODY_EN[0]
    return prs, slide


def save(prs: Presentation, out) -> None:
    prs.save(out)


# ── TRUE-POSITIVE vectors ─────────────────────────────────────────────────────
def notes(payload: str) -> Presentation:
    """Adversarial instruction buried in the speaker-notes pane."""
    prs, slide = _prs()
    slide.notes_slide.notes_text_frame.text = payload
    return prs


def off_slide(payload: str, *, x: int = -6_000_000, y: int = -6_000_000) -> Presentation:
    """Textbox placed off the slide canvas via negative EMU coordinates."""
    prs, slide = _prs()
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(3_000_000), Emu(600_000))
    tb.text_frame.text = payload
    return prs


def notes_multislide(payload: str, *, planted_on: int = 1, count: int = 3) -> Presentation:
    """Multi-slide deck; payload only in slide `planted_on` (0-based) notes."""
    prs = Presentation()
    for i in range(count):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = BODY_EN[i % len(BODY_EN)]
        if i == planted_on:
            slide.notes_slide.notes_text_frame.text = payload
    return prs


# ── TRUE-NEGATIVE / FP-hardening ──────────────────────────────────────────────
def plain() -> Presentation:
    prs, _ = _prs()
    return prs


def innocent_notes(text: str) -> Presentation:
    prs, slide = _prs()
    slide.notes_slide.notes_text_frame.text = text
    return prs


def visible_body(text: str) -> Presentation:
    """Text in a VISIBLE on-slide textbox (not notes, not off-slide). A fair
    warning belongs here — visible text is CLEAN even if it names an AI."""
    prs, slide = _prs()
    tb = slide.shapes.add_textbox(Emu(914_400), Emu(2_000_000), Emu(6_000_000), Emu(1_000_000))
    tb.text_frame.text = text
    return prs


# ── Robustness ────────────────────────────────────────────────────────────────
def corrupt_bytes() -> bytes:
    return b"PK\x03\x04" + bytes(range(256)) * 4


def encrypted_bytes() -> bytes:
    return bytes([0xD0, 0xCF, 0x11, 0xE0, 0xA1, 0xB1, 0x1A, 0xE1]) + bytes(64)
